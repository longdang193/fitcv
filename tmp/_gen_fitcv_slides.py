from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path('docs/slides/fitcv-overview.pptx')

# FitCV UI tokens (from src/fitcv_cp/templates/base.html, dark theme)
COLORS = {
    'bg': RGBColor(0x0F, 0x11, 0x17),
    'surface_1': RGBColor(0x1A, 0x1D, 0x2E),
    'surface_2': RGBColor(0x1E, 0x22, 0x35),
    'border': RGBColor(0x2D, 0x31, 0x48),
    'text_primary': RGBColor(0xE2, 0xE8, 0xF0),
    'text_secondary': RGBColor(0x94, 0xA3, 0xB8),
    'text_muted': RGBColor(0x64, 0x74, 0x8B),
    'accent': RGBColor(0x63, 0x66, 0xF1),
    'accent_hover': RGBColor(0x4F, 0x46, 0xE5),
    'accent_dim': RGBColor(0x31, 0x2E, 0x81),
    'accent_text': RGBColor(0xA5, 0xB4, 0xFC),
    'success': RGBColor(0x4A, 0xDE, 0x80),
    'error': RGBColor(0xF8, 0x71, 0x71),
    'warning': RGBColor(0xFB, 0x92, 0x3C),
    'info': RGBColor(0x60, 0xA5, 0xFA),
}

FONT = {
    'ui': 'Inter',
    'mono': 'Consolas',
}


def _no_line(shape) -> None:
    shape.line.fill.background()


def add_bg(slide, slide_w, slide_h) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg']
    _no_line(bg)


def add_nav(slide, slide_w, right_title: str | None = None) -> None:
    # Navigation bar mimic: height 56px-ish
    nav_h = Inches(0.55)
    nav = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, nav_h)
    nav.fill.solid()
    nav.fill.fore_color.rgb = COLORS['surface_1']
    nav.line.color.rgb = COLORS['border']
    nav.line.width = Pt(1)

    # Brand
    brand = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(3), Inches(0.3))
    tf = brand.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'FitCV'
    run.font.name = FONT['ui']
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = COLORS['accent']

    # Right-side label
    if right_title:
        box = slide.shapes.add_textbox(slide_w - Inches(4.2), Inches(0.15), Inches(3.6), Inches(0.3))
        tf2 = box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = right_title
        r2.font.name = FONT['ui']
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLORS['text_secondary']


def add_title(slide, text: str, x, y, w, h, size=36) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT['ui']
    r.font.bold = True
    r.font.size = Pt(size)
    r.font.color.rgb = COLORS['text_primary']


def add_subtitle(slide, text: str, x, y, w, h, size=16, color='text_secondary') -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT['ui']
    r.font.size = Pt(size)
    r.font.color.rgb = COLORS[color]


def add_card(slide, x, y, w, h, title: str | None = None, body: str | None = None) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS['surface_1']
    card.line.color.rgb = COLORS['border']
    card.line.width = Pt(1)

    pad_x = Inches(0.22)
    pad_y = Inches(0.18)
    if title:
        t = slide.shapes.add_textbox(x + pad_x, y + pad_y, w - 2 * pad_x, Inches(0.3))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT['ui']
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = COLORS['text_primary']
    if body:
        b = slide.shapes.add_textbox(x + pad_x, y + pad_y + Inches(0.38), w - 2 * pad_x, h - pad_y - Inches(0.38))
        tfb = b.text_frame
        tfb.clear()
        tfb.word_wrap = True
        p2 = tfb.paragraphs[0]
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT['ui']
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLORS['text_secondary']


def add_badge(slide, x, y, text: str, kind: str = 'info') -> None:
    bg_map = {
        'success': COLORS['accent_dim'],
        'info': COLORS['surface_2'],
        'warning': COLORS['surface_2'],
        'error': COLORS['surface_2'],
        'neutral': COLORS['surface_2'],
    }
    fg_map = {
        'success': COLORS['success'],
        'info': COLORS['info'],
        'warning': COLORS['warning'],
        'error': COLORS['error'],
        'neutral': COLORS['text_secondary'],
    }
    w = Inches(1.4)
    h = Inches(0.28)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_map.get(kind, COLORS['surface_2'])
    pill.line.color.rgb = COLORS['border']
    pill.line.width = Pt(1)

    t = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.03), w - Inches(0.24), h)
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT['ui']
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg_map.get(kind, COLORS['text_secondary'])


def add_button(slide, x, y, text: str) -> None:
    w = Inches(2.2)
    h = Inches(0.42)
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    btn.fill.solid()
    btn.fill.fore_color.rgb = COLORS['accent']
    _no_line(btn)
    t = slide.shapes.add_textbox(x, y + Inches(0.05), w, h)
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT['ui']
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_flow_step(slide, x, y, w, title, desc, accent=False) -> None:
    fill = COLORS['surface_1']
    border = COLORS['border']
    if accent:
        fill = COLORS['accent_dim']
        border = COLORS['accent']

    step = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.95))
    step.fill.solid()
    step.fill.fore_color.rgb = fill
    step.line.color.rgb = border
    step.line.width = Pt(1)

    t = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.13), w - Inches(0.4), Inches(0.3))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT['ui']
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = COLORS['text_primary']

    d = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.45), w - Inches(0.4), Inches(0.4))
    tfd = d.text_frame
    tfd.clear()
    p2 = tfd.paragraphs[0]
    r2 = p2.add_run()
    r2.text = desc
    r2.font.name = FONT['ui']
    r2.font.size = Pt(11)
    r2.font.color.rgb = COLORS['text_secondary']


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    s1 = prs.slides.add_slide(blank)
    add_bg(s1, slide_w, slide_h)
    add_nav(s1, slide_w)
    add_badge(s1, Inches(0.7), Inches(1.05), 'WHY FITCV', kind='neutral')
    add_title(s1, 'Evidence-first job matching\n+ CV generation', Inches(0.7), Inches(1.45), Inches(8.8), Inches(1.4), size=40)
    add_subtitle(
        s1,
        'Turn messy job posts into a reviewable shortlist — then generate tailored CV outputs\nonly when upstream evidence says “ready”.',
        Inches(0.7),
        Inches(3.05),
        Inches(10.5),
        Inches(0.7),
        size=15,
    )
    add_card(
        s1,
        Inches(0.7),
        Inches(4.1),
        Inches(5.95),
        Inches(2.6),
        title='What changes for job seekers',
        body='• Compare many roles with consistent fields\n• See why a role ranks higher (not vibes)\n• Stop rewriting — generate CV once fit is proven',
    )
    add_card(
        s1,
        Inches(6.9),
        Inches(4.1),
        Inches(5.75),
        Inches(2.6),
        title='Core idea',
        body='Evidence gates expensive work.\n\nFitCV keeps each run inspectable through artifacts and stage-owned truth, so you can iterate with confidence.',
    )
    add_button(s1, Inches(10.45), Inches(6.35), 'Start with evidence')

    # Slide 2 — Problem
    s2 = prs.slides.add_slide(blank)
    add_bg(s2, slide_w, slide_h)
    add_nav(s2, slide_w, 'Problem')
    add_title(s2, 'Job search feels like busywork', Inches(0.7), Inches(1.0), Inches(11.8), Inches(0.8), size=34)
    add_subtitle(s2, 'High volume + low signal turns applying into an exhausting, manual pipeline.', Inches(0.7), Inches(1.75), Inches(11.5), Inches(0.5))

    w = Inches(4.05)
    add_card(s2, Inches(0.7), Inches(2.55), w, Inches(1.55), title='Too many posts', body='You can’t read everything. You miss good roles or waste time on weak ones.')
    add_card(s2, Inches(4.65), Inches(2.55), w, Inches(1.55), title='Hard to compare', body='Every posting uses different language. Apples-to-apples ranking takes effort.')
    add_card(s2, Inches(8.6), Inches(2.55), w, Inches(1.55), title='CV tailoring tax', body='Small edits add up. Copy/paste quickly becomes hours per week.')

    add_card(
        s2,
        Inches(0.7),
        Inches(4.35),
        Inches(12.0),
        Inches(2.2),
        title='Result',
        body='You spend energy formatting and guessing — not improving fit, stories, and outcomes.',
    )

    # Slide 3 — Why existing workflows hurt
    s3 = prs.slides.add_slide(blank)
    add_bg(s3, slide_w, slide_h)
    add_nav(s3, slide_w, 'Why it hurts')
    add_title(s3, 'Existing workflows break at scale', Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.8), size=32)
    add_subtitle(s3, 'Manual pipelines create friction, inconsistency, and zero learning loop.', Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5))

    # Attached-tab card pattern
    outer = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.4), Inches(12.0), Inches(4.4))
    outer.fill.solid(); outer.fill.fore_color.rgb = COLORS['surface_1']
    outer.line.color.rgb = COLORS['border']; outer.line.width = Pt(1)

    tab = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.4), Inches(12.0), Inches(0.55))
    tab.fill.solid(); tab.fill.fore_color.rgb = COLORS['surface_1']
    tab.line.color.rgb = COLORS['border']; tab.line.width = Pt(1)

    def tab_btn(x, label, active=False):
        b = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.4), Inches(2.2), Inches(0.55))
        b.fill.solid(); b.fill.fore_color.rgb = COLORS['surface_1']
        b.line.fill.background()
        t = s3.shapes.add_textbox(x, Inches(2.52), Inches(2.2), Inches(0.35))
        tf=t.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=label; r.font.name=FONT['ui']; r.font.size=Pt(12)
        r.font.color.rgb = COLORS['accent'] if active else COLORS['text_secondary']
        r.font.bold = active
        if active:
            underline = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.45), Inches(2.92), Inches(1.3), Inches(0.04))
            underline.fill.solid(); underline.fill.fore_color.rgb = COLORS['accent']
            underline.line.fill.background()

    tab_btn(Inches(0.95), 'Copy/paste', active=True)
    tab_btn(Inches(3.15), 'Spreadsheets')
    tab_btn(Inches(5.35), 'Prompting')
    tab_btn(Inches(7.55), 'Version drift')

    add_card(s3, Inches(1.1), Inches(3.25), Inches(5.65), Inches(1.5), title='Pain', body='Inconsistent notes + no standard fields.\nSame role gets re-evaluated from scratch.')
    add_card(s3, Inches(6.95), Inches(3.25), Inches(5.45), Inches(1.5), title='Consequence', body='Hard to defend decisions.\nHard to learn what works from prior cycles.')
    add_badge(s3, Inches(1.15), Inches(5.05), 'low signal', kind='warning')
    add_badge(s3, Inches(2.7), Inches(5.05), 'rework', kind='error')
    add_badge(s3, Inches(4.25), Inches(5.05), 'no audit trail', kind='neutral')

    add_card(
        s3,
        Inches(1.1),
        Inches(5.45),
        Inches(11.3),
        Inches(1.1),
        title='What’s missing',
        body='A repeatable, inspectable workflow that turns messy inputs into stable decisions.',
    )

    # Slide 4 — How FitCV solves
    s4 = prs.slides.add_slide(blank)
    add_bg(s4, slide_w, slide_h)
    add_nav(s4, slide_w, 'Solution')
    add_title(s4, 'FitCV: evidence-first workflow', Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.8), size=32)
    add_subtitle(s4, 'Make decisions on structured evidence — then generate CV only when fit is proven.', Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5))

    add_flow_step(s4, Inches(0.7), Inches(2.6), Inches(3.9), '1) Ingest', 'Load many job posts\n(including noisy sources)')
    add_flow_step(s4, Inches(4.75), Inches(2.6), Inches(3.9), '2) Normalize + rank', 'Stable fields + explainable\nshortlist and ordering', accent=True)
    add_flow_step(s4, Inches(8.8), Inches(2.6), Inches(3.9), '3) Generate', 'CV outputs with validation\n+ repair safeguards')

    gate = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(3.7), Inches(3.9), Inches(0.55))
    gate.fill.solid(); gate.fill.fore_color.rgb = COLORS['surface_2']
    gate.line.color.rgb = COLORS['accent']; gate.line.width = Pt(1)
    t = s4.shapes.add_textbox(Inches(4.75), Inches(3.82), Inches(3.9), Inches(0.3))
    tf=t.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text='Only run “Generate” when READY'; r.font.name=FONT['ui']; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=COLORS['accent_text']

    add_card(
        s4,
        Inches(0.7),
        Inches(4.55),
        Inches(12.0),
        Inches(2.15),
        title='Why this matters',
        body='You stop paying “CV tailoring tax” up front. FitCV pushes effort downstream, after the shortlist is defensible.',
    )

    # Slide 5 — Key features
    s5 = prs.slides.add_slide(blank)
    add_bg(s5, slide_w, slide_h)
    add_nav(s5, slide_w, 'What you get')
    add_title(s5, 'Key capabilities (built for clarity)', Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.8), size=30)
    add_subtitle(s5, 'Features matter only if they remove pain and create a learning loop.', Inches(0.7), Inches(1.65), Inches(12.0), Inches(0.5))

    cw = Inches(6.0)
    ch = Inches(1.65)
    add_card(s5, Inches(0.7), Inches(2.45), cw, ch, title='Explainable ranking', body='See why roles score higher using consistent fields and evidence — not gut feel.')
    add_card(s5, Inches(7.0), Inches(2.45), cw, ch, title='Inspectable run artifacts', body='Every run keeps stage-owned truth so decisions can be reviewed, repeated, and improved.')
    add_card(s5, Inches(0.7), Inches(4.25), cw, ch, title='Control-plane UI', body='Component-like surfaces for reviewing inputs, adjusting settings, and confirming changes.')
    add_card(s5, Inches(7.0), Inches(4.25), cw, ch, title='Safe CV generation', body='Validation + repair safeguards reduce “almost good” drafts and wasted edits.')

    add_badge(s5, Inches(0.75), Inches(6.15), 'less rework', kind='success')
    add_badge(s5, Inches(2.25), Inches(6.15), 'more signal', kind='info')
    add_badge(s5, Inches(3.75), Inches(6.15), 'repeatable', kind='neutral')

    # Slide 6 — Why useful for job seekers
    s6 = prs.slides.add_slide(blank)
    add_bg(s6, slide_w, slide_h)
    add_nav(s6, slide_w, 'Why useful')
    add_title(s6, 'Why FitCV for job seekers', Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.8), size=32)
    add_subtitle(s6, 'Spend time where it increases outcomes: targeting, storytelling, and iteration.', Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5))

    add_card(s6, Inches(0.7), Inches(2.55), Inches(3.95), Inches(2.3), title='Fewer, better applications', body='Prioritize roles with the highest fit so effort goes to the best shots.')
    add_card(s6, Inches(4.85), Inches(2.55), Inches(3.95), Inches(2.3), title='Less rewriting', body='Generate CV output after fit is proven — stop tailoring for roles you won’t pursue.')
    add_card(s6, Inches(9.0), Inches(2.55), Inches(3.7), Inches(2.3), title='Clear decisions', body='Explainable ranking helps you defend “why this role” and refine your criteria.')

    add_card(
        s6,
        Inches(0.7),
        Inches(5.05),
        Inches(12.0),
        Inches(1.65),
        title='The loop',
        body='Import jobs → review shortlist → generate targeted CV → learn from outcomes → repeat with better signal.',
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == '__main__':
    build()
